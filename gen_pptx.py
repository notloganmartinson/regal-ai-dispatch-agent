from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

def create_pitch_deck():
    prs = Presentation()
    
    # J.B. Hunt Colors
    JB_YELLOW = RGBColor(255, 198, 39)
    DARK_BG = RGBColor(30, 30, 30)
    WHITE = RGBColor(255, 255, 255)

    slides_data = [
        {
            "title": "The Opportunity: J.B. Hunt AI Dispatch Triage",
            "content": [
                "Automating Fleet Crisis Routing and Policy Support",
                "The Challenge: Weather events and closures overwhelm human dispatchers, causing driver downtime.",
                "The Solution: A Level 1 Inbound Dispatch AI Agent for instant parallel-processing.",
                "The Impact: Handles routine FAQs and automated rerouting, reserving supervisors for unmapped escalations."
            ]
        },
        {
            "title": "How It Works: Logic, Knowledge, & Infrastructure",
            "content": [
                "Conversational Brain: Built on Vapi, tuned with explicit conversational flows for barge-ins.",
                "Knowledge Base (RAG): Custom FastAPI backend and PostgreSQL (pgvector) for semantic policy search.",
                "Escalation Protocol: Intelligent fallback hands off to a live supervisor if a reroute isn't found.",
                "Deployment: Hosted on Render to ensure 24/7 reliability and zero webhook latency."
            ]
        },
        {
            "title": "Driver Personalization: Real-Time SMS Routing",
            "content": [
                "Stateful Database Lookups: Dynamically references manifest, greeting by name and specific cargo.",
                "Cross-Channel Action: Pushes out-of-band communication directly to the driver's device.",
                "Technical Execution: Background task triggers Twilio API to send Google Maps reroutes.",
                "Note: Backend Twilio logic fires successfully; live SMS pending A2P 10DLC registration."
            ]
        },
        {
            "title": "Performance Data: 10-Call Simulation",
            "content": [
                "Intent Classification Accuracy: [XX]%",
                "Containment Rate: [XX]%",
                "Testing Breakdown:",
                "- [X] Standard Crisis Reroutes (Contained)",
                "- [X] Policy/FAQ Inquiries (Contained)",
                "- [X] Off-Topic Deflections (Contained)",
                "- [X] Live Handoffs (Escalated)"
            ]
        },
        {
            "title": "Product Vision: Next-Generation Dispatch",
            "content": [
                "Friction Point Identified: [Insert Placeholder: e.g., STT latency during background truck noise].",
                "Resilience & Queuing: Migrate Twilio SMS to Celery/Redis for automatic API retries.",
                "Proactive Telemetry: Integrate Google Maps Distance Matrix to text drivers before they reach closures."
            ]
        }
    ]

    for slide_data in slides_data:
        slide_layout = prs.slide_layouts[1] # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set dark background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

        # Format Title
        title = slide.shapes.title
        title.text = slide_data["title"]
        title.text_frame.paragraphs[0].font.color.rgb = JB_YELLOW
        title.text_frame.paragraphs[0].font.bold = True
        
        # Format Content
        content = slide.placeholders[1]
        content.text = slide_data["content"][0]
        content.text_frame.paragraphs[0].font.color.rgb = WHITE
        
        for point in slide_data["content"][1:]:
            p = content.text_frame.add_paragraph()
            p.text = point
            p.font.color.rgb = WHITE
            p.font.size = Pt(18)

    prs.save('Regal_AI_JBHunt_Pitch.pptx')
    print("Presentation saved successfully!")

if __name__ == '__main__':
    create_pitch_deck()
